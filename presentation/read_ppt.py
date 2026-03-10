from pptx import Presentation

prs = Presentation('Group no 10.pptx')

for i, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE {i+1}")
    print(f"{'='*60}")
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text.strip())
            print("---")
